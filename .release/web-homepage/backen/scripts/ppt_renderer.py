#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5


def load_json(path: str | None, fallback: Any) -> Any:
    if not path:
        return fallback
    p = Path(path)
    if not p.is_file():
        return fallback
    return json.loads(p.read_text(encoding="utf-8"))


def color(value: str, fallback: str = "1F2937") -> RGBColor:
    text = (value or fallback).strip().lstrip("#")
    if len(text) != 6:
        text = fallback
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def clean(value: Any, fallback: str = "") -> str:
    if value is None:
        return fallback
    text = str(value).strip()
    return text if text else fallback


def fit_text(value: str, limit: int) -> str:
    text = clean(value)
    return text if len(text) <= limit else text[: max(0, limit - 1)] + "…"


def remove_all_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx exposes no public delete API.
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def blank_layout(prs: Presentation):
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def image_index(images: list[str], manifest: dict[str, Any]) -> dict[str, str]:
    by_id: dict[str, str] = {}
    for image_path in images:
        path = Path(image_path)
        if path.is_file():
            by_id[path.stem] = str(path)
    for item in manifest.get("images", []) if isinstance(manifest, dict) else []:
        image_id = clean(item.get("id"))
        filename = clean(item.get("filename"))
        if not image_id:
            continue
        for image_path in images:
            path = Path(image_path)
            if path.stem == image_id or path.name == filename:
                by_id[image_id] = str(path)
                break
    return by_id


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int,
             font_color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    return box


def add_bullets(slide, bullets: list[Any], x: float, y: float, w: float, h: float, accent: RGBColor,
                font_color: RGBColor) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    values = [fit_text(item, 42) for item in bullets if clean(item)][:5]
    if not values:
        values = ["围绕主题提炼核心信息"]
    for index, item in enumerate(values):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = font_color
        p.space_after = Pt(8)
    marker = slide.shapes.add_shape(1, Inches(x - 0.18), Inches(y + 0.08), Inches(0.05), Inches(max(0.4, len(values) * 0.42)))
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent


def add_metrics(slide, metrics: list[Any], x: float, y: float, w: float, accent: RGBColor, text_color: RGBColor) -> None:
    values = [item for item in metrics if isinstance(item, dict)][:3]
    if not values:
        return
    card_w = w / len(values) - 0.12
    for index, item in enumerate(values):
        left = x + index * (card_w + 0.18)
        shape = slide.shapes.add_shape(1, Inches(left), Inches(y), Inches(card_w), Inches(0.92))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shape.line.color.rgb = accent
        add_text(slide, fit_text(item.get("value", ""), 14), left + 0.12, y + 0.12, card_w - 0.24, 0.3, 18, accent, True, PP_ALIGN.CENTER)
        add_text(slide, fit_text(item.get("label", ""), 18), left + 0.12, y + 0.48, card_w - 0.24, 0.25, 9, text_color, False, PP_ALIGN.CENTER)


def add_numbered_points(slide, bullets: list[Any], x: float, y: float, w: float, h: float, accent: RGBColor, text_color: RGBColor) -> None:
    values = [fit_text(item, 36) for item in bullets if clean(item)][:4]
    if not values:
        values = ["提炼关键依据"]
    row_h = min(0.78, h / max(len(values), 1))
    for index, item in enumerate(values):
        top = y + index * row_h
        badge = slide.shapes.add_shape(1, Inches(x), Inches(top + 0.05), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.color.rgb = accent
        add_text(slide, str(index + 1), x + 0.02, top + 0.105, 0.3, 0.18, 8, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        add_text(slide, item, x + 0.5, top, w - 0.55, 0.42, 15, text_color)


def render_metric_hero(slide, slide_data: dict[str, Any], accent: RGBColor, deep: RGBColor, text: RGBColor) -> None:
    metrics = [item for item in slide_data.get("metrics", []) if isinstance(item, dict)]
    hero = metrics[0] if metrics else {}
    value = fit_text(hero.get("value") or clean(slide_data.get("headline"), "核心发现"), 18)
    label = fit_text(hero.get("label") or clean(slide_data.get("title"), "关键结论"), 28)
    panel = slide.shapes.add_shape(1, Inches(0.95), Inches(2.05), Inches(4.35), Inches(2.15))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(248, 250, 252)
    panel.line.color.rgb = accent
    add_text(slide, value, 1.2, 2.34, 3.85, 0.65, 34, accent, True, PP_ALIGN.CENTER)
    add_text(slide, label, 1.2, 3.16, 3.85, 0.38, 13, deep, False, PP_ALIGN.CENTER)
    add_numbered_points(slide, slide_data.get("bullets", []), 6.0, 2.03, 5.85, 3.0, accent, text)


def render_comparison(slide, slide_data: dict[str, Any], accent: RGBColor, deep: RGBColor, text: RGBColor) -> None:
    bullets = [fit_text(item, 34) for item in slide_data.get("bullets", []) if clean(item)]
    midpoint = max(1, (len(bullets) + 1) // 2)
    groups = [bullets[:midpoint], bullets[midpoint:]]
    labels = ["现状 / 基线", "改进 / 发现"]
    for index, values in enumerate(groups):
        x = 0.95 + index * 5.85
        card = slide.shapes.add_shape(1, Inches(x), Inches(2.0), Inches(5.2), Inches(3.85))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        card.line.color.rgb = accent if index else RGBColor(203, 213, 225)
        add_text(slide, labels[index], x + 0.28, 2.24, 4.55, 0.3, 14, deep, True)
        add_bullets(slide, values or ["保留核心对照信息"], x + 0.45, 2.83, 4.25, 2.35, accent, text)


def render_matrix(slide, slide_data: dict[str, Any], accent: RGBColor, deep: RGBColor, text: RGBColor) -> None:
    bullets = [fit_text(item, 28) for item in slide_data.get("bullets", []) if clean(item)][:4]
    while len(bullets) < 4:
        bullets.append("补充分析维度")
    for index, item in enumerate(bullets):
        row, col = divmod(index, 2)
        x = 1.0 + col * 5.7
        y = 1.95 + row * 1.9
        tile = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.05), Inches(1.38))
        tile.fill.solid()
        tile.fill.fore_color.rgb = RGBColor(248, 250, 252)
        tile.line.color.rgb = RGBColor(226, 232, 240)
        add_text(slide, f"{index + 1:02d}", x + 0.22, y + 0.22, 0.55, 0.25, 12, accent, True)
        add_text(slide, item, x + 0.95, y + 0.23, 3.75, 0.55, 15, deep if index == 0 else text)


def render_timeline(slide, slide_data: dict[str, Any], accent: RGBColor, text: RGBColor) -> None:
    bullets = [fit_text(item, 28) for item in slide_data.get("bullets", []) if clean(item)][:5]
    if not bullets:
        bullets = ["问题定义", "方法构建", "实验验证", "结论输出"]
    y = 3.24
    line = slide.shapes.add_shape(1, Inches(1.05), Inches(y + 0.16), Inches(10.9), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    step = 10.6 / max(len(bullets) - 1, 1)
    for index, item in enumerate(bullets):
        x = 1.0 + index * step
        dot = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.34), Inches(0.34))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.color.rgb = accent
        add_text(slide, item, x - 0.42, y + 0.55, 1.45, 0.78, 10, text, False, PP_ALIGN.CENTER)


def image_box(image_path: str, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    path = Path(image_path)
    try:
        with Image.open(path) as image:
            iw, ih = image.size
    except Exception:
        return x, y, w, h
    if iw <= 0 or ih <= 0:
        return x, y, w, h
    scale = min(w / iw, h / ih)
    rendered_w = iw * scale
    rendered_h = ih * scale
    return x + (w - rendered_w) / 2, y + (h - rendered_h) / 2, rendered_w, rendered_h


def add_image_in_box(slide, image_path: str, x: float, y: float, w: float, h: float, accent: RGBColor | None = None) -> None:
    path = Path(image_path)
    if not path.is_file():
        return
    if accent:
        frame = slide.shapes.add_shape(1, Inches(x - 0.04), Inches(y - 0.04), Inches(w + 0.08), Inches(h + 0.08))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
        frame.line.color.rgb = accent
    rx, ry, rw, rh = image_box(str(path), x, y, w, h)
    slide.shapes.add_picture(str(path), Inches(rx), Inches(ry), width=Inches(rw), height=Inches(rh))


def add_image(slide, image_path: str, layout: str, accent: RGBColor | None = None) -> None:
    if layout == "full-image":
        add_image_in_box(slide, image_path, 0.78, 1.25, 11.78, 5.48, accent)
    elif layout == "image-left":
        add_image_in_box(slide, image_path, 0.75, 1.55, 5.45, 4.85, accent)
    elif layout == "image-top":
        add_image_in_box(slide, image_path, 1.0, 1.44, 11.25, 3.25, accent)
    elif layout == "evidence-strip":
        add_image_in_box(slide, image_path, 0.82, 2.0, 4.2, 3.9, accent)
    else:
        add_image_in_box(slide, image_path, 7.0, 1.55, 5.45, 4.85, accent)


def add_footer(slide, deck: dict[str, Any], index: int, total: int, muted: RGBColor) -> None:
    add_text(slide, fit_text(deck.get("title", "AI 生成 PPT"), 40), 0.72, 7.08, 8.0, 0.2, 7, muted)
    add_text(slide, f"{index + 1:02d}/{total:02d}", 11.9, 7.05, 0.8, 0.24, 8, muted, False, PP_ALIGN.RIGHT)


def render_cover(slide, deck: dict[str, Any], slide_data: dict[str, Any], palette: list[str]) -> None:
    accent = color(palette[0])
    deep = color(palette[1])
    pale = color(palette[3], "EFF6FF")
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = pale
    rail = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.28), Inches(SLIDE_H))
    rail.fill.solid()
    rail.fill.fore_color.rgb = accent
    add_text(slide, fit_text(slide_data.get("title") or deck.get("title") or "AI 生成 PPT", 34), 0.9, 1.6, 10.6, 1.1, 34, deep, True)
    add_text(slide, fit_text(slide_data.get("headline") or deck.get("subtitle") or deck.get("audience") or "", 62), 0.94, 2.9, 9.7, 0.62, 17, color(palette[4]))
    add_text(slide, fit_text(deck.get("theme") or "Generated presentation", 42), 0.96, 5.72, 5.8, 0.32, 11, color(palette[4]))


def render_section(slide, deck: dict[str, Any], slide_data: dict[str, Any], palette: list[str], index: int, total: int) -> None:
    accent = color(palette[0])
    deep = color(palette[1])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = deep
    add_text(slide, clean(slide_data.get("section"), "SECTION"), 0.82, 1.38, 3.6, 0.35, 13, color(palette[2], "D9A441"), True)
    add_text(slide, fit_text(slide_data.get("title"), 24), 0.82, 2.22, 9.8, 0.92, 34, RGBColor(255, 255, 255), True)
    add_text(slide, fit_text(slide_data.get("headline"), 60), 0.86, 3.34, 8.4, 0.48, 15, RGBColor(226, 232, 240))
    rule = slide.shapes.add_shape(1, Inches(0.86), Inches(4.36), Inches(1.8), Inches(0.06))
    rule.fill.solid()
    rule.fill.fore_color.rgb = accent
    add_footer(slide, deck, index, total, RGBColor(203, 213, 225))


def render_content(slide, deck: dict[str, Any], slide_data: dict[str, Any], palette: list[str],
                   index: int, total: int, image_path: str | None) -> None:
    accent = color(palette[0])
    deep = color(palette[1])
    text = color(palette[4])
    layout = clean(slide_data.get("layout"), "auto")
    add_text(slide, clean(slide_data.get("section"), clean(slide_data.get("type"), "SLIDE")).upper(), 0.72, 0.48, 3.4, 0.26, 9, accent, True)
    add_text(slide, fit_text(slide_data.get("title"), 34), 0.72, 0.82, 8.4, 0.55, 24, deep, True)
    if clean(slide_data.get("headline")):
        add_text(slide, fit_text(slide_data.get("headline"), 64), 0.74, 1.34, 9.2, 0.36, 12, text)
    if layout == "metric-hero":
        render_metric_hero(slide, slide_data, accent, deep, text)
    elif layout == "comparison":
        render_comparison(slide, slide_data, accent, deep, text)
        add_metrics(slide, slide_data.get("metrics", []), 1.0, 6.02, 11.0, accent, text)
    elif layout == "matrix":
        render_matrix(slide, slide_data, accent, deep, text)
    elif layout == "timeline":
        render_timeline(slide, slide_data, accent, text)
    elif image_path and layout == "full-image":
        add_image(slide, image_path, "full-image", accent)
    elif image_path and layout == "image-top":
        add_image(slide, image_path, "image-top", accent)
        add_numbered_points(slide, slide_data.get("bullets", []), 1.18, 5.05, 10.8, 1.05, accent, text)
    elif image_path and layout == "evidence-strip":
        add_image(slide, image_path, "evidence-strip", accent)
        add_numbered_points(slide, slide_data.get("bullets", []), 5.55, 2.05, 6.2, 3.2, accent, text)
        add_metrics(slide, slide_data.get("metrics", []), 5.55, 5.75, 6.1, accent, text)
    elif image_path and layout == "image-left":
        add_image(slide, image_path, "image-left", accent)
        add_bullets(slide, slide_data.get("bullets", []), 6.8, 2.0, 5.6, 3.65, accent, text)
        add_metrics(slide, slide_data.get("metrics", []), 6.8, 5.8, 5.4, accent, text)
    elif image_path:
        add_bullets(slide, slide_data.get("bullets", []), 0.95, 2.0, 5.5, 3.65, accent, text)
        add_metrics(slide, slide_data.get("metrics", []), 0.95, 5.8, 5.4, accent, text)
        add_image(slide, image_path, "image-right", accent)
    else:
        add_bullets(slide, slide_data.get("bullets", []), 1.05, 2.0, 10.8, 3.55, accent, text)
        add_metrics(slide, slide_data.get("metrics", []), 1.05, 5.75, 10.8, accent, text)
    add_footer(slide, deck, index, total, RGBColor(148, 163, 184))


def render_deck(args: argparse.Namespace) -> None:
    deck = load_json(args.deck, {})
    style = load_json(args.style, {})
    images = load_json(args.images, [])
    manifest = load_json(args.manifest, {"images": []})
    template = Path(args.template) if args.template else None

    prs = Presentation(str(template)) if template and template.is_file() else Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    remove_all_slides(prs)

    palette = style.get("palette") if isinstance(style.get("palette"), list) else ["005BAC", "063A78", "D9A441", "EFF6FF", "1F2937"]
    palette = (palette + ["005BAC", "063A78", "D9A441", "EFF6FF", "1F2937"])[:5]
    by_id = image_index(images if isinstance(images, list) else [], manifest if isinstance(manifest, dict) else {})
    slides = deck.get("slides", []) if isinstance(deck, dict) else []
    if not slides:
        raise ValueError("deck.slides is empty")
    if len(slides) > 40:
        raise ValueError(f"deck.slides exceeds 40-slide limit: {len(slides)}")

    total = len(slides)
    for index, item in enumerate(slides):
        slide_data = item if isinstance(item, dict) else {}
        slide = prs.slides.add_slide(blank_layout(prs))
        slide_type = clean(slide_data.get("type"), "content")
        image_path = by_id.get(clean(slide_data.get("imageId")))
        if index == 0 or slide_type == "cover":
            render_cover(slide, deck, slide_data, palette)
        elif slide_type == "section":
            render_section(slide, deck, slide_data, palette, index, total)
        else:
            render_content(slide, deck, slide_data, palette, index, total, image_path)

    output = Path(args.out)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(json.dumps({"ok": True, "out": str(output), "slides": len(slides)}, ensure_ascii=False))


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--deck", required=True)
    parser.add_argument("--style", required=True)
    parser.add_argument("--images", required=True)
    parser.add_argument("--manifest", required=False)
    parser.add_argument("--template", required=False)
    parser.add_argument("--out", required=True)
    render_deck(parser.parse_args())
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
